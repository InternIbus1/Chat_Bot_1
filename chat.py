import streamlit as st
import os
import time
import io
import re
import base64
import tempfile
import pandas as pd
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from PIL import Image
import google.generativeai as genai
import requests
from bs4 import BeautifulSoup
from datetime import datetime
from html import unescape
import json
import pickle
import base64
from streamlit.components.v1 import html



# --- Constants ---
CHAT_CACHE_DIR = "chat_cache"
EXTS = ['pdf', 'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls', 'csv']
IMAGE_EXTS = ['jpg', 'jpeg', 'png', 'gif', 'bmp']

# Ensure chat cache directory exists
if not os.path.exists(CHAT_CACHE_DIR):
    os.makedirs(CHAT_CACHE_DIR)

# Set page config - MUST BE THE FIRST STREAMLIT COMMAND
st.set_page_config(
    "iBUS Chatbot",
    "💬", # This is the emoji for the tab icon
    "wide",
    initial_sidebar_state="expanded",
    menu_items={
        'Get Help': 'https://www.ibusnetworks.com/help',
        'Report a bug': 'https://www.ibusnetworks.com/bug',
        'About': 'iBUS Networks Interactive Chatbot'
    }
)

# --- CSS Styling ---
st.markdown("""
<style>
    /* Override Streamlit's default header behavior */
    header {
        visibility: hidden;
    }
    
    /* Global styling for components inside sidebar for consistency */
    /* Targeting a known Streamlit sidebar container class */
    .st-emotion-cache-1pxazr7 { 
        background-color: white;
        padding: 20px; /* Add some padding */
    }

    /* Style for Streamlit native image component in sidebar */
    /* Target img directly within the sidebar container */
    .st-emotion-cache-1pxazr7 img { 
        border-radius: 8px;
        margin-bottom: 10px;
        max-width: 70px;
        height: auto;
        display: block;
        margin-left: auto;
        margin-right: auto;
    }

    /* Adjust Streamlit h1 in sidebar */
    /* Target h1 directly within the sidebar container's markdown elements */
    .st-emotion-cache-1pxazr7 h1 { 
        font-size: 24px;
        color: #003A6C;
        margin-bottom: 5px;
        text-align: center; /* Center the title */
    }
    
    /* Adjust Streamlit text in sidebar */
    /* Target p directly within the sidebar container's markdown elements */
    .st-emotion-cache-1pxazr7 p { 
        font-size: 14px;
        color: #6C757D;
        margin-bottom: 20px;
        text-align: center; /* Center the description */
    }

    /* Separator in sidebar */
    .st-emotion-cache-1pxazr7 hr {
        border-top: 1px solid #E9ECEF;
        margin-top: 20px;
        margin-bottom: 20px;
    }

</style>
""", unsafe_allow_html=True)


# Main content wrapper (no longer needs margin-top as header is moved)
st.markdown("""
<div class="main-content-wrapper">
""", unsafe_allow_html=True)

st.markdown("""
<style>
    /* Enhanced chat container with iBUS branding */
    .chat-container {
        display: flex;
        flex-direction: column;
        gap: 12px;
        margin-bottom: 18px;
        width: 100%;
    }
    
    /* User message - right aligned with iBUS primary color */
    .user-message-container {
        display: flex;
        justify-content: flex-end;
        width: 100%;
    }
    
    .user-message {
        background-color: var(--ibus-primary);
        color: white;
        border-radius: 18px 18px 0 18px;
        padding: 12px 18px;
        max-width: 80%;
        margin-left: auto;
        box-shadow: 0 2px 6px rgba(0,0,0,0.15);
        position: relative;
    }
    
    /* Bot message - left aligned with light background */
    .bot-message-container {
        display: flex;
        justify-content: flex-start;
        width: 100%;
    }
    
    .bot-message {
        background-color: white;
        color: var(--ibus-primary);
        border-radius: 18px 18px 18px 0;
        padding: 12px 18px;
        max-width: 80%;
        margin-right: auto;
        box-shadow: 0 2px 6px rgba(0,0,0,0.1);
        position: relative;
        border-left: 3px solid var(--ibus-secondary);
    }
    
    /* Username labels with iBUS colors */
    .username-label {
        font-size: 0.85em;
        margin-bottom: 6px;
        font-weight: 600;
    }
    
    .user-label {
        text-align: right;
        color: var(--ibus-primary);
    }
    
    .bot-label {
        text-align: left;
        color: var(--ibus-secondary);
    }
    
    /* Improved timestamp styling */
    .timestamp {
        font-size: 0.7em;
        color: rgba(255, 255, 255, 0.9);
        margin-top: 5px;
        text-align: right;
        display: inline-block;
        float: right;
        clear: both;
        width: 100%;
    }
    
    .bot-message .timestamp {
        color: rgba(0, 58, 108, 0.7);
    }
    
    /* Message content with better spacing */
    .message-content {
        display: inline-block;
        width: 100%;
        margin-bottom: 8px;
        line-height: 1.4;
    }
    
    /* Enhanced animation for new messages */
    @keyframes slideInRight {
        from { opacity: 0; transform: translateX(20px); }
        to { opacity: 1; transform: translateX(0); }
    }
    
    @keyframes slideInLeft {
        from { opacity: 0; transform: translateX(-20px); }
        to { opacity: 1; transform: translateX(0); }
    }
    
    .user-message {
        animation: slideInRight 0.3s ease-out;
    }
    
    .bot-message {
        animation: slideInLeft 0.3s ease-out;
    }
    
    /* Add subtle hover effect */
    .user-message:hover, .bot-message:hover {
        box-shadow: 0 4px 8px rgba(0,0,0,0.2);
    }
</style>
""", unsafe_allow_html=True)

# Enhanced CSS styling based on iBUS logo colors (blue and gray tones)
st.markdown("""
<style>
    /* Color palette based on iBUS logo */
    :root {
        --ibus-primary: #003A6C;     /* Dark blue from logo */
        --ibus-secondary: #0077B6;   /* Medium blue */
        --ibus-accent: #48CAE4;      /* Light blue accent */
        --ibus-light: #ADE8F4;       /* Very light blue */
        --ibus-gray: #6C757D;        /* Complementary gray */
        --ibus-light-gray: #F8F9FA;  /* Background gray */
    }
    
    /* Global styling */
    .stApp {
        background-color: #F8F9FA;
    }
    
    /* Header styling */
    h1, h2, h3 {
        color: var(--ibus-primary);
        font-family: 'Helvetica Neue', Arial, sans-serif;
    }
    
    h1 {
        font-weight: 600;
        letter-spacing: -0.5px;
    }
    
    /* Sidebar styling */
    .css-1d391kg, .css-12oz5g7 { /* These are often the main sidebar and main app containers */
        background-color: white;
        border-right: 1px solid #E9ECEF;
    }
    
    /* Button styling */
    .stButton button {
        background-color: var(--ibus-primary);
        color: white;
        border-radius: 8px;
        border: none;
        padding: 10px 15px;
        font-weight: 500;
        transition: all 0.3s ease;
    }
    
    .stButton button:hover {
        background-color: var(--ibus-secondary);
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    
    /* Primary button */
    .stButton button[data-baseweb="button"][kind="primary"] {
        background-color: var(--ibus-primary);
    }
    
    /* Secondary button */
    .stButton button[data-baseweb="button"][kind="secondary"] {
        background-color: var(--ibus-gray);
    }
    
    /* Chat message styling */
    .stChatMessage {
        border-radius: 12px;
        padding: 10px;
        margin: 10px 0;
        box-shadow: 0 2px 5px rgba(0,0,0,0.05);
    }
    
    /* User message */
    .stChatMessage[data-testid="stChatMessage-USER"] {
        background-color: var(--ibus-light);
        border-bottom-right-radius: 4px;
    }
    
    /* Assistant message */
    .stChatMessage[data-testid="stChatMessage-ASSISTANT"] {
        background-color: white;
        border-bottom-left-radius: 4px;
    }
    
    /* Input box styling */
    .stTextInput input {
        border-radius: 8px;
        border: 1px solid #CED4DA;
        padding: 12px;
        box-shadow: 0 2px 5px rgba(0,0,0,0.05);
        transition: all 0.3s ease;
    }
    
    .stTextInput input:focus {
        border-color: var(--ibus-secondary);
        box-shadow: 0 0 0 3px rgba(0,119,182,0.2);
    }
    
    /* File uploader styling */
    .stFileUploader {
        background-color: white;
        border-radius: 8px;
        padding: 10px;
        border: 1px dashed #CED4DA;
    }
    
    /* Expander styling */
    .streamlit-expanderHeader {
        background-color: white;
        border-radius: 8px;
        border: none;
        box-shadow: 0 2px 5px rgba(0,0,0,0.05);
        padding: 10px 15px;
        font-weight: 500;
        color: var(--ibus-primary);
    }
    
    .streamlit-expanderContent {
        background-color: white;
        border-radius: 0 0 8px 8px;
        padding: 15px;
        border: none;
        box-shadow: 0 2px 5px rgba(0,0,0,0.05);
    }
    
    /* Logo container (not used in fixed header anymore) */
    .logo-container {
        display: flex;
        align-items: center;
        padding: 15px;
        background-color: white;
        border-radius: 12px;
        box-shadow: 0 4px 10px rgba(0,0,0,0.05);
        margin-bottom: 20px;
    }
    
    .logo-container img {
        height: 60px;
        margin-right: 15px;
    }
    
    .logo-container h1 {
        margin: 0;
        color: var(--ibus-primary);
    }
    
    /* Fade-in animation for elements */
    @keyframes fadeIn {
        from { opacity: 0; transform: translateY(10px); }
        to { opacity: 1; transform: translateY(0); }
    }
    
    /* Apply animations to different components */
    .stButton button {
        transition: all 0.3s ease;
    }
    
    /* Chat message animations */
    .element-container:has(.stChatMessage) {
        animation: fadeIn 0.5s ease-out forwards;
    }
    
    /* Header animations */
    h1, h2, h3 {
        animation: fadeIn 0.7s ease-out forwards;
    }
    
    /* Staggered animation delays */
    .staggered-1 { animation-delay: 0.2s; }
    .staggered-2 { animation-delay: 0.4s; }
    .staggered-3 { animation-delay: 0.6s; }
    .staggered-4 { animation-delay: 0.8s; }
    
    /* Welcome message special animation */
    .welcome-message {
        animation: fadeIn 0.8s ease-out forwards;
    }
    
    /* Options container animation */
    .options-container {
        animation: fadeIn 1s ease-out forwards;
        animation-delay: 0.5s;
    }
    
    /* Typing effect */
    .typing-effect {
        border-left: 2px solid var(--ibus-secondary);
        padding-left: 8px;
        color: var(--ibus-primary);
        font-weight: 500;
    }
    
    /* Card styling for options */
    .option-card {
        background-color: white;
        border-radius: 8px;
        padding: 15px;
        box-shadow: 0 4px 8px rgba(0,0,0,0.05);
        transition: all 0.3s ease;
        border-left: 4px solid var(--ibus-primary);
        margin: 10px 0;
        cursor: pointer;
    }
    
    .option-card:hover {
        transform: translateY(-3px);
        box-shadow: 0 6px 12px rgba(0,0,0,0.1);
        border-left: 44px solid var(--ibus-secondary);
    }
    
    /* Follow-up suggestions styling */
    .followup-container {
        margin-top: 20px;
        padding: 15px;
        background-color: var(--ibus-light-gray);
        border-radius: 8px;
        border-left: 4px solid var(--ibus-accent);
    }

    /* Make follow-up buttons horizontal */
    .followup-container + div > .st-emotion-cache-nahz7x { /* Targeting the element containing columns for follow-up buttons */
        display: flex;
        flex-wrap: wrap; /* Allow wrapping to next line if space is limited */
        gap: 10px; /* Space between buttons */
        justify-content: flex-start; /* Align buttons to the start */
    }

    .followup-container + div > .st-emotion-cache-nahz7x .stButton button { 
        flex-grow: 1; /* Allow buttons to grow and fill available space */
        min-width: 150px; /* Minimum width for each button */
        max-width: 300px; /* Maximum width to prevent them from becoming too wide */
        white-space: normal !important; /* Allow text to wrap within the button */
        height: auto !important; /* Allow height to adjust based on content */
        display: inline-block !important; /* Ensure it behaves well in flexbox */
        text-align: center !important; /* Center text within button */
        padding: 10px 15px !important; /* Adjust padding */
        font-size: 0.9rem !important; /* Smaller font size */
        margin: 0 !important; /* Remove default margin */
    }
    
    /* Success message styling */
    .success-message {
        padding: 15px;
        background-color: #d4edda;
        color: #155724;
        border-radius: 8px;
        margin: 15px 0;
        border-left: 4px solid #28a745;
        animation: fadeIn 0.5s ease-out forwards;
    }
    
</style>
""", unsafe_allow_html=True)

st.markdown("""
<style>
    /* Improved typing indicator with iBUS colors */
    .typing-indicator {
        display: flex;
        align-items: center;
        justify-content: flex-start;
        height: 24px;
        padding-left: 10px;
    }
    
    .typing-indicator span {
        height: 8px;
        width: 8px;
        background-color: var(--ibus-secondary);
        border-radius: 50%;
        display: inline-block;
        margin: 0 3px;
        opacity: 0.4;
    }
    
    .typing-indicator span:nth-child(1) {
        animation: pulse 1s infinite;
    }
    
    .typing-indicator span:nth-child(2) {
        animation: pulse 1s infinite 0.2s;
    }
    
    .typing-indicator span:nth-child(3) {
        animation: pulse 1s infinite 0.4s;
    }
    
    @keyframes pulse {
        0% { opacity: 0.4; transform: scale(1); }
        50% { opacity: 1; transform: scale(1.2); }
        100% { opacity: 0.4; transform: scale(1); }
    }
</style>
""", unsafe_allow_html=True)

st.markdown("""
<style>
    /* Timestamp styling - make it more robust */
    .timestamp {
        font-size: 0.7em;
        color: black;
        margin-top: 5px;
        text-align: right;
        display: inline-block;
        float: right;
        clear: both;
        width: 100%;
    }
    
    /* Ensure message content and timestamp don't overlap */
    .message-content {
        display: inline-block;
        width: 100%;
        margin-bottom: 15px;
    }
</style>
""", unsafe_allow_html=True)


st.markdown("""
<style>
    /* File attachment button styling */
    .attachment-button {
        background-color: transparent;
        border: none;
        color: var(--ibus-primary); /* Use iBUS primary color for the icon */
        font-size: 24px; /* Larger icon */
        cursor: pointer;
        padding: 8px; /* Padding for click area */
        border-radius: 50%;
        transition: background-color 0.2s, transform 0.2s;
        display: flex; /* To center icon */
        align-items: center;
        justify-content: center;
    }
    
    .attachment-button:hover {
        background-color: var(--ibus-light); /* Light blue hover */
        transform: scale(1.1);
    }
    
    /* iBUS Mascot styling - fixed positioning */
    .ibus-mascot {
        position: absolute;
        left: 10px;
        bottom: 10px;
        z-index: 100;
        display: flex;
        align-items: center;
        justify-content: center;
        width: 30px;
        height: 30px;
        overflow: visible;
    }
    
    .ibus-mascot img {
        width: 30px;
        height: 30px;
        border-radius: 50%;
        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
        object-fit: cover;
    }
    
    /* File upload area styling */
    .file-upload-area {
        margin-top: 10px; /* Space between input and uploader */
        margin-bottom: 10px;
        padding: 10px;
        border-radius: 8px;
        background-color: #f7f7f7;
        border: 1px dashed #CED4DA;
    }
    /* Style for the file uploader widget itself */
    .stFileUploader {
        background-color: white;
        border-radius: 8px;
        padding: 10px;
        border: 1px dashed #CED4DA;
    }

    /* Scroll buttons styling */
    .scroll-buttons {
        position: fixed;
        bottom: 90px; /* Above the chat input */
        right: 20px;
        display: flex;
        flex-direction: column;
        gap: 10px;
        z-index: 999; /* Ensure they are above other content */
    }

    .scroll-button {
        background-color: var(--ibus-secondary);
        color: white;
        border: none;
        border-radius: 50%;
        width: 40px;
        height: 40px;
        display: flex;
        align-items: center;
        justify-content: center;
        font-size: 20px;
        cursor: pointer;
        box-shadow: 0 2px 5px rgba(0,0,0,0.2);
        transition: background-color 0.3s, transform 0.2s;
    }

    .scroll-button:hover {
        background-color: var(--ibus-primary);
        transform: scale(1.1);
    }
</style>
""", unsafe_allow_html=True)

st.markdown("""
<style>
    /* User avatar styling */
    .user-avatar {
        position: absolute;
        right: 10px;
        bottom: 10px;
        z-index: 100;
        display: flex;
        align-items: center;
        justify-content: center;
        width: 30px;
        height: 30px;
        border-radius: 50%;
        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
        object-fit: cover;
    }
    
    .user-avatar img {
        width: 30px;
        height: 30px;
        border-radius: 50%;
        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
        object-fit: cover;
    }
    
    /* Update user message to make room for avatar */
    .user-message {
        position: relative;
        padding-right: 40px; /* Add space for avatar */
    }
</style>
""", unsafe_allow_html=True)


# Configure Gemini
genai.configure(api_key="AIzaSyAc2xYCF0uCE6Ti3sTjc9-wxyvI5hHtBSM") # !!! REPLACE WITH YOUR ACTUAL API KEY !!!

# --- Chat History Functions ---
def save_chat_history(user_name):
    """Save chat history and other session data to disk."""
    if user_name:
        cache_file = os.path.join(CHAT_CACHE_DIR, f'{user_name.lower().replace(" ", "_")}.pkl')
        user_data = {
            'chat_history': st.session_state.chat_history,
            'message_timestamps': st.session_state.message_timestamps,
            'documents_content': st.session_state.documents_content,
            'processed_files': st.session_state.processed_files,
            'file_summaries': st.session_state.file_summaries,
            'images': st.session_state.images,
            'user_name': st.session_state.user_name,
            'asked_name': st.session_state.asked_name,
            'show_file_upload_expander': st.session_state.show_file_upload_expander,
            'files_displayed': st.session_state.files_displayed,
            'show_options': st.session_state.get('show_options', False),
            'current_followups': st.session_state.get('current_followups', []),
        }
        try:
            with open(cache_file, 'wb') as f:
                pickle.dump(user_data, f)
            # st.success("Chat history saved!") # Can be removed for silent saving
        except Exception as e:
            st.error(f"Error saving chat history: {e}")

def load_chat_history(user_name):
    """Load chat history and other session data from disk."""
    if user_name:
        cache_file = os.path.join(CHAT_CACHE_DIR, f'{user_name.lower().replace(" ", "_")}.pkl')
        if os.path.exists(cache_file):
            try:
                with open(cache_file, 'rb') as f:
                    user_data = pickle.load(f)

                st.session_state.chat_history = user_data.get('chat_history', [])
                st.session_state.message_timestamps = user_data.get('message_timestamps', {})
                st.session_state.documents_content = user_data.get('documents_content', {})
                st.session_state.processed_files = user_data.get('processed_files', [])
                st.session_state.images = user_data.get('images', [])
                
                # Only load user_name if it matches the current user trying to load
                if user_data.get('user_name') == user_name:
                    st.session_state.user_name = user_data.get('user_name')
                    st.session_state.asked_name = user_data.get('asked_name', False)
                    st.session_state.show_file_upload_expander = user_data.get('show_file_upload_expander', False)
                    st.session_state.files_displayed = user_data.get('files_displayed', False)
                    st.session_state.show_options = user_data.get('show_options', False)
                    st.session_state.current_followups = user_data.get('current_followups', [])

                    # Ensure file_summaries exists in session state before assigning
                    if 'file_summaries' not in st.session_state:
                        st.session_state['file_summaries'] = {}
                    st.session_state['file_summaries'].update(user_data.get('file_summaries', {}))
                    return True
                else:
                    # If user name doesn't match, treat as new session for this user
                    return False
            except Exception as e:
                st.error(f"Error loading chat history for {user_name}: {e}")
    return False

def get_base64_image(image_path):
    with open(image_path, "rb") as f:
        data = f.read()
    return base64.b64encode(data).decode()

# Check if logo.jpg exists, otherwise provide a placeholder
if os.path.exists("logo.jpg"):
    image_base64 = get_base64_image("logo.jpg")
    image_tag = f'<img src="data:image/jpeg;base64,{image_base64}" alt="iBUS Mascot" width="30" height="30">'
else:
    image_tag = '<span style="font-size: 30px;">🤖</span>' # Fallback emoji


# --- Text Extraction ---
def extract_text(p, ext):
    try:
        if ext == 'pdf': return "".join(pg.get_text() for pg in fitz.open(p))
        if ext in ('docx', 'doc'): return "\n".join(par.text for par in docx.Document(p).paragraphs)
        if ext in ('pptx', 'ppt'): return "\n".join(s.text for sl in Presentation(p).slides for s in sl.shapes if hasattr(s, 'text'))
        if ext == 'csv':
            df = pd.read_csv(p)
            return df.to_string()
        if ext in ('xlsx', 'xls'):
            df = pd.read_excel(p)
            return df.to_string()
    except Exception as e: st.error(f"{os.path.basename(p)}: {e}")

# --- User Avatar Generation ---
def generate_user_avatar(username):
    """Generate a user avatar with the first letter of their name"""
    if not username:
        return ""
    
    # Get first letter and capitalize it
    first_letter = username[0].upper()
    
    # Generate a consistent color based on the username
    hash_value = sum(ord(c) for c in username)
    hue = hash_value % 360  # 0-359 degrees on color wheel
    
    # Create a vibrant but not too light color (HSL format)
    bg_color = f"hsl({hue}, 70%, 60%)"
    
    # Create the SVG avatar
    svg = f'''
    <svg width="30" height="30" viewBox="0 0 30 30" xmlns="http://www.w3.org/2000/svg">
        <circle cx="15" cy="15" r="15" fill="{bg_color}"/>
        <text x="15" y="20" font-family="Arial, sans-serif" font-size="16" 
              font-weight="bold" fill="white" text-anchor="middle">{first_letter}</text>
    </svg>
    '''
    
    # Return the SVG as a data URI
    return f'<img src="data:image/svg+xml;base64,{base64.b64encode(svg.encode()).decode()}" alt="User Avatar">'


# --- Gemini Interaction ---
def ask_gemini(q, ctx, images=None):
    try:
        mdl = genai.GenerativeModel('gemini-1.5-flash')
        
        # Check if the question is asking for real-time or current information
        real_time_indicators = [
            "current", "latest", "today", "now", "recent", "update", 
            "real-time", "real time", "live", "news", "weather", "stock", 
            "price", "market", "trending", "happening", "this week",
            "this month", "this year", "forecast", "prediction"
        ]
        
        # Detect if this is a general knowledge question or iBUS specific
        ibus_indicators = ["ibus", "network", "telecom", "company", "service", "infrastructure"]
        is_ibus_related = any(indicator in q.lower() for indicator in ibus_indicators)
        
        # Check if it needs real-time data
        needs_real_time = any(indicator in q.lower() for indicator in real_time_indicators)
        
        # For general knowledge questions that aren't about iBUS, we can ignore the context
        if not is_ibus_related and ("what is" in q.lower() or "how to" in q.lower() or "who is" in q.lower()):
            # This is likely a general knowledge question
            ctx = "You are a helpful assistant that can answer general knowledge questions."
        
        # Ensure we have context for iBUS-related questions
        if is_ibus_related and (not ctx or ctx == "iBUS Networks is a telecommunications company."):
            # Check if we have document content in session state
            if st.session_state.documents_content:
                ctx = "\n".join(f"{n}:\n{c}" for n, c in st.session_state.documents_content.items())
        
        # Prepare prompt with context and explicit instruction not to repeat the question
        base_prompt = f"""Context:
{ctx}

Question: {q}

Important: Provide a direct answer without repeating the question. Do not use formats like "Question: ... Answer: ...". Just give the answer. If the answer is not directly available in the context, use your knowledge to provide the most accurate and helpful response. Never say that you don't have information or that the answer isn't in the provided documents - always try to give a helpful response."""
        
        # If real-time data is needed, add a web search component
        if needs_real_time:
            try:
                # Add current date and time information
                current_time_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                real_time_info = f"Current date and time: {current_time_str}\n\n"
                
                # Attempt to get some basic real-time information
                search_terms = q.replace("?", "").replace("!", "").replace(".", "")
                search_url = f"https://news.google.com/rss/search?q={search_terms}&hl=en-US&gl=US&ceid=US:en"
                
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
                }
                
                response = requests.get(search_url, headers=headers, timeout=5)
                # Parse RSS feed to extract relevant snippets (simplified for brevity)
                soup = BeautifulSoup(response.content, 'xml')
                items = soup.find_all('item')
                snippets = []
                for item in items[:3]: # Get top 3 news items
                    title = unescape(item.find('title').text)
                    link = item.find('link').text
                    snippets.append(f"Title: {title}\nLink: {link}\n")
                
                if snippets:
                    real_time_info += "Recent News Snippets:\n" + "\n---\n".join(snippets) + "\n\n"

                # Add the real-time information to the prompt with the same instruction
                prompt = f"""{base_prompt}

Real-time information:
{real_time_info}

Please provide an up-to-date answer based on both the context and the real-time information. If the question is about general knowledge and not related to iBUS Networks, focus on the real-time information. Remember, do not repeat the question in your answer and never say you don't have information."""
            except Exception as e:
                # If web search fails, fall back to the base prompt with a note
                prompt = f"{base_prompt}\n\nNote: I attempted to search for real-time information but encountered an error: {str(e)}. I'll answer based on my available knowledge."
        else:
            prompt = base_prompt
        
        if images:
            parts = [{"text": prompt}] + [
                {"inline_data": {"mime_type": "image/jpeg", "data": img}} for img in images
            ]
            response = mdl.generate_content(contents=[{"parts": parts}]).text
        else:
            response = mdl.generate_content(contents=[{"parts":[{"text": prompt}]}]).text
        
        # Clean up the response to remove any question repetition and leading/trailing asterisks
        if "Question:" in response and "Answer:" in response:
            response = response.split("Answer:", 1)[1].strip()
        
        # Remove any leading/trailing asterisks that might cause formatting issues
        response = response.strip('*')
        
        return response
    except Exception as e: 
        # Check if it's a rate limit error
        error_str = str(e)
        if "429" in error_str or "quota" in error_str.lower() or "rate limit" in error_str.lower():
            return """I'm currently experiencing high demand and have reached my API rate limit. Please try again in a few minutes, or consider:
1. Asking simpler questions
2. Breaking your request into smaller parts
3. Using fewer images in your queries

This is a temporary limitation of the free tier API usage."""
        else:
            return f"Error: {e}"

# Generate follow-up questions based on the conversation
def generate_followups(q, a, ctx):
    try:
        # Create a more detailed prompt for generating specific follow-up questions
        prompt = f"""Based on this conversation:
User: {q}
Assistant: {a}

Generate 3 specific, detailed follow-up questions the user might want to ask next. Each question should:
1. Be directly related to the topic discussed
2. Ask for more specific details or examples
3. Explore a natural next step in the conversation
4. Be phrased as a complete question (with a question mark)
5. Be concise (under 10 words if possible)
6. Avoid using the word "you" - phrase questions objectively

Format as a simple numbered list (1. 2. 3.) with no additional text."""
        
        try:
            response = ask_gemini(prompt, ctx)
            # Check if we got a rate limit error response
            if "rate limit" in response.lower() or "quota" in response.lower():
                raise Exception("Rate limit exceeded")
            
            # --- IMPROVED PARSING OF FOLLOW-UP QUESTIONS ---
            questions = []
            for line in response.split('\n'):
                line = line.strip()
                # Check if the line starts with a number followed by a period and optional space
                if re.match(r'^\d+\.\s*', line):
                    # Remove the number and period, then strip
                    q_text = re.sub(r'^\d+\.\s*', '', line).strip()
                    if q_text:
                        questions.append(q_text)
            # --- END IMPROVED PARSING ---
        
            specific_questions = []
            for question in questions:
                # Add a question mark if missing
                if not question.endswith('?'):
                    question += '?'
                # Ensure the question is not too generic by adding context-specific terms
                # This part is still good but refined with more robust keyword extraction
                
                # Extract potential keywords from the full conversation context (q, a, ctx)
                full_text_for_keywords = f"{q} {a} {ctx}"
                words = re.findall(r'\b\w+\b', full_text_for_keywords.lower())
                
                # Filter out common words and get potential key terms (e.g., nouns, proper nouns)
                common_words = {'the', 'and', 'is', 'in', 'to', 'a', 'of', 'for', 'with', 'on', 'at', 'what', 'how', 'who', 'when', 'where', 'why', 'can', 'you', 'i', 'me', 'it', 'its', 'this', 'that', 'these', 'those', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'not', 'no', 'yes', 'or', 'but', 'if', 'so', 'as', 'by', 'from', 'up', 'down', 'out', 'in', 'on', 'off', 'about', 'just', 'more', 'some', 'any', 'all', 'etc'}
                
                # Prioritize capitalized words first, then longer words
                potential_keywords = sorted(
                    list(set([w for w in words if w not in common_words and len(w) > 3])),
                    key=lambda x: (x[0].isupper(), len(x)), reverse=True
                )
                
                # Ensure the question is specific
                generic_patterns = [
                    r'tell me more', r'can you explain', r'what else', r'how does', r'why is',
                    r'further details', r'more information', r'elaborate on', r'details about',
                    r'explain this', r'what about'
                ]
                
                is_generic = any(re.search(pattern, question.lower()) for pattern in generic_patterns)
                
                if is_generic and potential_keywords:
                    for keyword in potential_keywords:
                        if keyword.lower() not in question.lower():
                            # Append keyword if it makes sense, e.g., "Tell me more about X"
                            if re.search(r'(tell me more|explain|elaborate)( on)?( this)?\??', question.lower()):
                                question = re.sub(r'(tell me more|explain|elaborate)( on)?( this)?\??', f'What about {keyword}?', question.lower(), 1)
                            else: # For other generic questions, just append
                                question = f"{question.rstrip('?')} about {keyword}?"
                            break # Use only one keyword to keep it concise
                
                specific_questions.append(question.strip().replace('  ', ' ')) # Clean up extra spaces
                
                if len(specific_questions) >= 3:
                    break # Stop after getting 3 good questions

            # If we don't have enough specific questions from Gemini's output
            if len(specific_questions) < 3:
                # Use a more sophisticated fallback based on the original question's keywords
                q_words = q.lower().split()
                q_keywords = [w for w in q_words if w not in common_words and len(w) > 3]
            
                fallback_options = []
                if q_keywords:
                    # Try to generate questions related to the core topic of the user's last input
                    fallback_options.append(f"What are the benefits of {q_keywords[0]}?")
                    if len(q_keywords) > 1:
                        fallback_options.append(f"How does {q_keywords[0]} relate to {q_keywords[1]}?")
                    fallback_options.append(f"Can you provide an example of {q_keywords[0]}?")
                else:
                    # Default generic fallback if no good keywords from the question
                    fallback_options = [
                        "What specific features does this offer?",
                        "How is this implemented in practice?",
                        "What are the next steps I should take?"
                    ]
                
                # Fill remaining spots with fallback options
                for fo in fallback_options:
                    if len(specific_questions) < 3 and fo not in specific_questions:
                        specific_questions.append(fo)
            
            return specific_questions[:3] # Ensure exactly 3 questions
        
        except Exception as e:
            # Fallback for API errors or other generation issues
            print(f"Error in follow-up generation: {e}")
            # Return more specific generic questions if there's an error
            return [
                "What are the key benefits of this approach?",
                "How does this compare to alternatives?",
                "Can you provide a specific example." # Changed from '?' to '.' to ensure no issues with the last char.
            ]
    except Exception as e:
        print(f"Outer error generating follow-ups: {e}")
        return [
            "What are the key benefits of this approach?",
            "How does this compare to alternatives?",
            "Can you provide a specific example." # Changed from '?' to '.' to ensure no issues with the last char.
        ]

# --- UI Helper Functions ---
def show_loading_animation(seconds=1.5):
    """Display a loading animation for the specified number of seconds"""
    progress_placeholder = st.empty()
    progress_bar = progress_placeholder.progress(0)
    for i in range(100):
        time.sleep(seconds / 100)
        progress_bar.progress(i + 1)
    progress_placeholder.empty()

# def handle_predefined_option(option):
#     """Handle clicks on predefined option buttons"""
#     if option == "Upload Files":
#         st.session_state['show_file_upload_expander'] = not st.session_state['show_file_upload_expander'] # Toggle visibility
#         return # Do not add to chat history as it's an action

#     elif option == "What is iBUS?":
#         message = "iBUS Networks is a leading telecommunications company specializing in innovative connectivity solutions for businesses and organizations."
#     elif option == "Services offered":
#         message = "iBUS Networks offers a range of services including:\n\n- High-speed internet connectivity\n- Network infrastructure solutions\n- Cloud services\n- Managed IT services\n- Telecommunications consulting"
#     elif option == "Contact information":
#         message = "You can contact iBUS Networks through:\n\n- Email: info@ibusnetworks.com\n- Phone: +1-555-IBUS-NET\n- Website: www.ibusnetworks.com"
#     elif option == "Help with this chatbot":
#         message = "This chatbot can help you with:\n\n1. Information about iBUS Networks and services\n2. Analyzing documents you upload\n3. Answering questions about telecommunications\n\nJust type your question or upload files to get started!"
#     elif option == "Predict":
#         st.write("Click the button below to open the Rainfall Prediction Dashboard:")
#         if st.button("噫 Open Rainfall Prediction Dashboard"):
#             js = '''
#             <script>
#             window.open("https://rainfall-prediction-dashboard.streamlit.app/", "_blank");
#             </script>
#             '''
#             st.components.v1.html(js)
#     elif option == "Site Location Map":
#         # Automatically open the link in a new tab
#             js = """
#             <script>
#                 window.open("http://localhost:3000", "_blank");
#             </script>
#             """
#             st.markdown(js, unsafe_allow_html=True)
#             message = "Opening Site Map......."
#     else:
#         message = f"You selected: {option}"

#     # Add the user's selection to chat history
#     st.session_state.chat_history.append({"role": "user", "content": option})
#     # Add the response to chat history
#     st.session_state.chat_history.append({"role": "assistant", "content": message})
#     # Generate follow-up questions based on this interaction
#     ctx = "iBUS Networks is a telecommunications company."
#     st.session_state['current_followups'] = generate_followups(option, message, ctx)
#     # Turn off the options display after selection
#     st.session_state['show_options'] = False
#     save_chat_history(st.session_state.user_name) # Save history after every interaction
#     st.rerun()


def handle_predefined_option(option):
    """Handle clicks on predefined option buttons"""
    message = None # Initialize message to None to prevent UnboundLocalError

    if option == "Upload Files":
        st.session_state['show_file_upload_expander'] = not st.session_state['show_file_upload_expander'] # Toggle visibility
        return # Do not add to chat history as it's an action
    elif option == "What is iBUS?":
        message = "iBUS Networks is a leading telecommunications company specializing in innovative connectivity solutions for businesses and organizations."
    elif option == "Services offered":
        message = "iBUS Networks offers a range of services including:\n\n- High-speed internet connectivity\n- Network infrastructure solutions\n- Cloud services\n- Managed IT services\n- Telecommunications consulting"
    elif option == "Contact information":
        message = "You can contact iBUS Networks through:\n\n- Email: info@ibusnetworks.com\n- Phone: +1-555-IBUS-NET\n- Website: www.ibusnetworks.com"
    elif option == "Help with this chatbot":
        message = "This chatbot can help you with:\n\n1. Information about iBUS Networks and services\n2. Analyzing documents you upload\n3. Answering questions about telecommunications\n\nJust type your question or upload files to get started!"
    elif option == "Predict":
        message="Opening the rainfall prediction model.."
        js = "window.open('https://rainfall-prediction-dashboard.streamlit.app/')"
        st.markdown(js, unsafe_allow_html=True)
        st.components.v1.html(f"<script>{js}</script>")
           
            # A message for the chat history is already set above.
    elif option == "Site Location Map":
        # Assign a message that will go into the chat history for this action
        message = "Opening Site Location Map in a new tab......."
        # Automatically open the link in a new tab
        js = "window.open('https://sitemap1.netlify.app/')"
        st.markdown(js, unsafe_allow_html=True)
        st.components.v1.html(f"<script>{js}</script>")
    else:
        # Fallback for any unhandled predefined options
        message = f"You selected: {option}"
    
    # Only add to chat history if a message was successfully generated
    if message is not None:
        st.session_state.chat_history.append({"role": "user", "content": option, "id": f"user_{len(st.session_state.chat_history)}"})
        assistant_time = datetime.now().strftime("%I:%M %p")
        message_id = f"assistant_{len(st.session_state.chat_history)}"
        st.session_state.message_timestamps[message_id] = assistant_time
        st.session_state.chat_history.append({"role": "assistant", "content": message, "id": message_id})
        save_chat_history(st.session_state.user_name)
        st.rerun() # Use parentheses for function call

def new_chat_session():
    """Clears the current chat history and starts a new session."""
    if st.session_state.user_name:
        save_chat_history(st.session_state.user_name) # Save current chat before clearing
    
    st.session_state.chat_history = []
    st.session_state.documents_content = {}
    st.session_state.processed_files = []
    st.session_state.images = []
    st.session_state.file_summaries = {}
    st.session_state.message_timestamps = {}
    st.session_state.files_displayed = False
    st.session_state.show_file_upload_expander = False # Reset expander visibility
    st.session_state.show_options = True # Show options for new chat
    st.session_state.current_followups = []
    st.session_state.temp_user_input = None # Clear any pending temp input

# --- Session State Initialization (MUST BE FIRST) ---
# Initialize all session state variables at the very beginning
for k, v in {'chat_history': [], 'documents_content': {}, 'processed_files': [], 'images': [],
             'user_name': None, 'asked_name': False, 'file_summaries': {}, 
             'show_file_upload_expander': False, 'message_timestamps': {}, 
             'files_displayed': False, 'show_options': True, 'current_followups': [],
             'temp_user_input': None}.items(): # Added temp_user_input here
    st.session_state.setdefault(k, v)


# --- Main Application Logic ---

# Sidebar for new chat
with st.sidebar:
    # Moved header content into sidebar using native Streamlit components
    st.image("logo.jpg" if os.path.exists("logo.jpg") else "https://via.placeholder.com/70x70?text=Logo", width=70) 
    st.title("iBUS Interactive Chatbot")
    st.write("Your intelligent telecommunications assistant")
    st.markdown("<hr>", unsafe_allow_html=True) # Separator
    
    st.title("Chat Options")
    col1, col2 = st.columns(2)
    with col1:
        if st.button("New Chat", on_click=new_chat_session):
            pass # The function handles rerun
    with col2:
        if st.button("Clear All", on_click=new_chat_session): # Reuse new_chat_session to clear all
            pass # The function handles rerun


# Ask for name if not already asked
if not st.session_state.user_name:
    if not st.session_state.asked_name:
        current_time = datetime.now().strftime("%I:%M %p")
        greeting_msg = "Hello! I'm the iBUS chatbot. Can I know your name, please?"
        
        # Display the greeting with a typing effect
        greeting_placeholder = st.empty()
        time.sleep(0.8) # Add a slight delay before showing the greeting
        for i in range(1, len(greeting_msg) + 1):
            greeting_placeholder.markdown(f"<div class='typing-effect'>{greeting_msg[:i]}</div>", unsafe_allow_html=True)
            time.sleep(0.03) # Adjust typing speed
        greeting_placeholder.markdown(f"""
            <div class="chat-container">
                <div class="username-label bot-label">iChat</div>
                <div class="bot-message-container">
                    <div class="bot-message">
                        <div class="message-content">{greeting_msg}</div>
                        <div class="timestamp">{current_time}</div>
                        <div class="ibus-mascot">
                           {image_tag}
                        </div>
                    </div>
                </div>
            </div>
            """, unsafe_allow_html=True)
        # Store the greeting message in chat history with a placeholder ID
        message_id = f"assistant_greeting"
        st.session_state.message_timestamps[message_id] = current_time
        st.session_state.chat_history.append({"role": "assistant", "content": greeting_msg, "id": message_id})
        st.session_state.asked_name = True
        # No need to save history here, it will be saved after name input
        
    user_name_input = st.text_input("", placeholder="Enter your name here...", key="name_input")
    if user_name_input:
        st.session_state.user_name = user_name_input
        # Check if history exists for this user
        if load_chat_history(st.session_state.user_name):
            st.success(f"Welcome back, {st.session_state.user_name}! Loading your chat history.")
        else:
            # New user, add their name to history and a welcome message
            user_time = datetime.now().strftime("%I:%M %p")
            message_id = f"user_{len(st.session_state.chat_history)}"
            st.session_state.message_timestamps[message_id] = user_time
            st.session_state.chat_history.append({"role": "user", "content": user_name_input, "id": message_id})

            welcome_time = datetime.now().strftime("%I:%M %p")
            welcome_msg = f"👋 Welcome {user_name_input}! How can I assist you today?"
            message_id = f"assistant_{len(st.session_state.chat_history)}"
            st.session_state.message_timestamps[message_id] = welcome_time
            st.session_state.chat_history.append({"role": "assistant", "content": welcome_msg, "id": message_id})
            st.session_state['show_options'] = True
            st.session_state['show_staggered_animation'] = True
        save_chat_history(st.session_state.user_name) # Save history after name input/load
        st.rerun()
else:
    # Display full chat history if we already have the user's name
    for i, m in enumerate(st.session_state.chat_history):
        message_id = m.get('id', f"{m['role']}_{i}")
        timestamp = st.session_state.message_timestamps.get(message_id, "")
        if m['role'] == 'user':
            st.markdown(f"""
                <div class="chat-container">
                    <div class="username-label user-label">{st.session_state.user_name}</div>
                    <div class="user-message-container">
                        <div class="user-message">
                            <div class="message-content">{m['content']}</div>
                            <div class="timestamp">{timestamp}</div>
                            <div class="user-avatar">
                                {generate_user_avatar(st.session_state.user_name)}
                            </div>
                        </div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
        else:
            content = m['content']
            if "Question:" in content and "Answer:" in content:
                content = content.split("Answer:", 1)[1].strip()
            content = content.lstrip('*') # Remove any leading asterisks
            st.markdown(f"""
                <div class="chat-container">
                    <div class="username-label bot-label">iChat</div>
                    <div class="bot-message-container">
                        <div class="bot-message">
                            <div class="message-content">{content}</div>
                            <div class="timestamp">{timestamp}</div>
                            <div class="ibus-mascot">
                                {image_tag}
                            </div>
                        </div>
                    </div>
                </div>
                """, unsafe_allow_html=True)

    # Show predefined options after welcome message if flag is set
    if st.session_state.get('show_options', False) and not st.session_state.processed_files:
        time.sleep(0.5) # Add a slight delay before showing options
        st.markdown("""
            <div class="options-container">
                <h3 style="color: #003A6C; margin-bottom: 15px;"> <span style="margin-right: 8px;">💡</span>How can I help you today?</h3>
            </div>
            """, unsafe_allow_html=True)
        
        # Display buttons for predefined options with staggered animation
        options = ["What is iBUS?", "Services offered", "Contact information", "Upload Files", "Help with this chatbot","Predict","Site Location Map"]
        for i, option in enumerate(options):
            st.button(
                option,
                on_click=handle_predefined_option,
                args=(option,),
                key=f"option_{i}",
                help=f"Click to learn about {option}",
            )

    # Place the chat input and file uploader at the bottom
    st.markdown("<div style='height: 100px;'></div>", unsafe_allow_html=True) # Empty space to push chat input down

    # File Upload Expander - shown or hidden based on session state, moved ABOVE chat input
    if st.session_state.get('show_file_upload_expander'):
        with st.expander("Upload Documents", expanded=True):
            uploaded_files = st.file_uploader(
                "Upload documents for analysis (PDF, DOCX, PPTX, CSV, XLSX, Images)",
                type=EXTS + IMAGE_EXTS,
                accept_multiple_files=True,
                key="inline_file_uploader"
            )
            if uploaded_files:
                for uploaded_file in uploaded_files:
                    file_name = uploaded_file.name
                    file_extension = file_name.split('.')[-1].lower()

                    if file_name not in st.session_state.processed_files: # Prevent re-processing
                        if file_extension in EXTS:
                            try:
                                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_extension}") as tmp_file:
                                    tmp_file.write(uploaded_file.getvalue())
                                temp_path = tmp_file.name
                                
                                text_content = extract_text(temp_path, file_extension)
                                if text_content:
                                    st.session_state.documents_content[file_name] = text_content
                                st.session_state.processed_files.append(file_name)
                                st.success(f"Processed {file_name} successfully.")
                                os.unlink(temp_path)
                            except Exception as e:
                                st.error(f"Error processing {file_name}: {e}")
                        elif file_extension in IMAGE_EXTS:
                            try:
                                image_bytes = uploaded_file.getvalue()
                                st.session_state.images.append(base64.b64encode(image_bytes).decode('utf-8'))
                                st.session_state.processed_files.append(file_name)
                                st.success(f"Processed image {file_name} successfully.")
                            except Exception as e:
                                st.error(f"Error processing image {file_name}: {e}")
                
                # After processing, add a message to chat history about the files
                if st.session_state.processed_files:
                    file_list = ", ".join(st.session_state.processed_files)
                    file_upload_message = f"I've processed the following files for context: {file_list}. You can now ask questions about them!"
                    
                    user_time = datetime.now().strftime("%I:%M %p")
                    message_id = f"user_file_upload_{len(st.session_state.chat_history)}"
                    st.session_state.message_timestamps[message_id] = user_time
                    st.session_state.chat_history.append({"role": "user", "content": f"Uploaded files: {file_list}", "id": message_id})

                    assistant_time = datetime.now().strftime("%I:%M %p")
                    message_id = f"assistant_file_ack_{len(st.session_state.chat_history)}"
                    st.session_state.message_timestamps[message_id] = assistant_time
                    st.session_state.chat_history.append({"role": "assistant", "content": file_upload_message, "id": message_id})
                    
                    # Generate follow-up questions relevant to the uploaded files
                    files_context = "\n".join(st.session_state.documents_content.values())
                    st.session_state['current_followups'] = generate_followups(f"Summarize these documents: {file_list}", file_upload_message, files_context)

                    st.session_state['show_file_upload_expander'] = False # Hide expander after successful upload and message
                    save_chat_history(st.session_state.user_name)
                    st.rerun()

    # Display follow-up questions if available, moved ABOVE chat input
    if st.session_state.get('current_followups'):
        st.markdown("""
            <div class="followup-container">
                <h4 style="color: var(--ibus-primary); margin-bottom: 10px;">Suggestions:</h4>
            </div>
            """, unsafe_allow_html=True)
        # Use a single container to apply flexbox styling for horizontal buttons
        st.container()
        cols = st.columns(len(st.session_state['current_followups'])) # Still use columns for semantic grouping
        for i, followup_q in enumerate(st.session_state['current_followups']):
            with cols[i]:
                # On click, directly add to history and rerun
                if st.button(followup_q, key=f"followup_{i}"):
                    st.session_state.temp_user_input = followup_q # Set temporary input
                    st.session_state['current_followups'] = [] # Clear follow-ups after selection
                    save_chat_history(st.session_state.user_name) # Save history after follow-up click
                    st.rerun() # Rerun to process the new user input

    # Create a container for the chat input and attachment button
    chat_input_container = st.container()
    with chat_input_container:
        # Create columns for the chat input and the attachment button
        # Adjusted column widths slightly and removed extra markdown for better alignment
        cols = st.columns([0.9, 0.1])
        with cols[0]:
            # Removed 'value' parameter as it caused TypeError
            user_input = st.chat_input("Ask me anything...", key="chat_input")
        with cols[1]:
            # Simple direct CSS for alignment - might need more fine-tuning depending on Streamlit's internal rendering
            st.markdown("<div style='margin-top: 14px;'>", unsafe_allow_html=True) # Adjust margin-top to align vertically
            if st.button("📎", key="attach_file_button", help="Attach files for context"):
                st.session_state['show_file_upload_expander'] = not st.session_state['show_file_upload_expander']
                st.rerun() # Added rerun here to immediately reflect the expander's visibility
            st.markdown("</div>", unsafe_allow_html=True)

    # Use the temporary input if available, otherwise use regular chat input
    input_to_process = user_input
    if st.session_state.temp_user_input:
        input_to_process = st.session_state.temp_user_input
        st.session_state.temp_user_input = None # Clear after use

    if input_to_process and st.session_state.user_name:
        current_time = datetime.now().strftime("%I:%M %p")
        message_id = f"user_{len(st.session_state.chat_history)}"
        st.session_state.message_timestamps[message_id] = current_time
        # Only append if it's new input (not already appended by followup button)
        if not any(msg.get('content') == input_to_process and msg.get('role') == 'user' for msg in st.session_state.chat_history[-2:]):
             st.session_state.chat_history.append({"role": "user", "content": input_to_process, "id": message_id})
        
        # Clear options after user types
        st.session_state['show_options'] = False
        st.session_state['current_followups'] = [] # Clear previous follow-ups

        with st.spinner("Thinking..."):
            show_loading_animation(2) # Show loading animation for 2 seconds

            ctx = "iBUS Networks is a telecommunications company."
            if st.session_state.documents_content:
                ctx += "\n\n" + "\n".join(f"{name}:\n{content}" for name, content in st.session_state.documents_content.items())
            
            response = ask_gemini(input_to_process, ctx, images=st.session_state.images)
            
            assistant_time = datetime.now().strftime("%I:%M %p")
            message_id = f"assistant_{len(st.session_state.chat_history)}"
            st.session_state.message_timestamps[message_id] = assistant_time
            st.session_state.chat_history.append({"role": "assistant", "content": response, "id": message_id})
            
            # Generate follow-up questions
            st.session_state['current_followups'] = generate_followups(input_to_process, response, ctx)
        
        save_chat_history(st.session_state.user_name) # Save history after every interaction
        st.rerun()

# End of main content wrapper
st.markdown("</div>", unsafe_allow_html=True)



